'''
@Author: tao_jingfen
@Date: 2020-06-30 13:41:44
@LastEditors: tao_jingfen
@LastEditTime: 2020-07-03 10:21:09
'''

from pptx import Presentation
import math
from pptx.util import Cm
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

def purificationStepPageMake(prs: str, stepsDict: dict) -> str:
    """
    Create the every step protein purification info page of purification auto-report PPT
    
    --Args:
        prs: presentation object loaded from *pptx* by template pptx
        stepsDict: every step protein purification info Dict,
        e.g. {
        1: {'PurificationStep': 'Protein A & Ni',
        'WXXX-hPro1.His - P1': 
        {
            'No':'1',
            'Concentration': 'XXX',
            'Volume': 'XXX',
            'Amount': 'XXX',
            'Yield': 'XXX',
            'Buffer': 'XXX',
            'Purity by SEC-HPLC': 'XXX',
            'Supernatant':'XXX',
            'MW':'XXX',
            'PI': 'XXX',
            'Comments': 'XXX'
        },
        'WXXX-hPro1.His - P2': 
        {
            'No':'1',
            'Concentration': 'XXX',
            'Volume': 'XXX',
            'Amount': 'XXX',
            'Yield': 'XXX',
            'Buffer': 'XXX',
            'Purity by SEC-HPLC': 'XXX',
            'Supernatant':'XXX',
            'MW':'XXX',
            'PI': 'XXX',
            'Comments': 'XXX'
        },
        'WXXX-BMK1': 
        {
            'No':'2',
            'Concentration': 'XXX',
            'Volume': 'XXX',
            'Amount': 'XXX',
            'Yield': 'XXX',
            'Buffer': 'XXX',
            'Purity by SEC-HPLC': 'XXX',
            'Supernatant':'XXX',
            'MW':'XXX',
            'PI': 'XXX',
            'Comments': 'XXX'
        }
        },
        2: {'PurificationStep': 'Dialysis',
        'WXXX-hPro1.His - P1': {},
        'WXXX-hPro1.His - P2': {},
        'WXXX-BMK1': {} }
        ...}
    --Returns:
        prs?presentation object could be used to save pptx
        2 purification steps with 5  each page at most
        more will generate another page
    """
    site = 1
#    if stepsDict == {}:
#        prs = stepTableMake(prs,1,'Protein A',{},site)
#    else:
    for n in stepsDict.keys():
        step = stepsDict[n].get('PurificationStep','')
        stepInfo = stepsDict.get(n,{})
        stepSampleNum = len(stepInfo.keys())-1
        if site == 1:
            prs = stepTableMake(prs,n,step,stepInfo,site)
            if stepSampleNum <= 5:
                site = 2
        else:
            if stepSampleNum <= 5:
                prs = stepTableMake(prs,n,step,stepInfo,site)
                site = 1
            else:
                site = 1
                prs = stepTableMake(prs,n,step,stepInfo,site)
    return prs
        
def stepTableMake(prs,n,step,stepInfo,site):
    """
    Create the protein purification info table
    
    --Args:
        prs: presentation object loaded from *pptx* by template pptx
        n: the number of the step, e.g. 1, 2, 3, start from 1
        step: the step name, e.g. 'Protein A & Ni'
        stepInfo: the step detailed info dict 
        e.g.
        {'PurificationStep': ,
        'WXXX-hPro1.His - P1': 
        {
            'No':'1',
            'Concentration': 'XXX',
            'Volume': 'XXX',
            'Amount': 'XXX',
            'Yield': 'XXX',
            'Buffer': 'XXX',
            'Purity by SEC-HPLC': 'XXX',
            'Supernatant':'XXX',
            'MW':'XXX',
            'PI': 'XXX',
            'Comments': 'XXX'
        },
        'WXXX-hPro1.His - P2': 
        {
            'No':'1',
            'Concentration': 'XXX',
            'Volume': 'XXX',
            'Amount': 'XXX',
            'Yield': 'XXX',
            'Buffer': 'XXX',
            'Purity by SEC-HPLC': 'XXX',
            'Supernatant':'XXX',
            'MW':'XXX',
            'PI': 'XXX',
            'Comments': 'XXX'
        },
        'WXXX-BMK1': 
        {
            'No':'2',
            'Concentration': 'XXX',
            'Volume': 'XXX',
            'Amount': 'XXX',
            'Yield': 'XXX',
            'Buffer': 'XXX',
            'Purity by SEC-HPLC': 'XXX',
            'Supernatant':'XXX',
            'MW':'XXX',
            'PI': 'XXX',
            'Comments': 'XXX'
        }
        },
        2: {'PurificationStep': 'Dialysis',
        'WXXX-hPro1.His - P1': {},
        'WXXX-hPro1.His - P2': {},
        'WXXX-BMK1': {} }
        ...}
        site: the site that place the table, 1 or 2
    --Returns:
        prs: presentation object could be used to save pptx
    """
    if site == 1:
        slide= prs.slides.add_slide(prs.slide_layouts[2])
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = 'Characterization'
        subtitle = title_shape.text_frame.add_paragraph()
        subtitle.text = 'Steps '
        subtitle.font.italic = True 
        subtitle1 = subtitle.add_run()
        subtitle1.text = '- Protein Purification Summary'
        subtitle1.font.bold = False
        subtitle1.font.italic = True
        top = Cm(2.8)
    else:
        slide = prs.slides[-1]
        shapes = slide.shapes
        top = Cm(10.8)
    shape = shapes.add_textbox(Cm(2.2), top, Cm(8.65), Cm(1.11))
    if 'Protein A' in step or 'Ni' in step:
        column = 11
    else:
        column = 12
    table = shapes.add_table(max(1,len(stepInfo.keys())), column, Cm(0.6), top + Cm(1.11), Cm(34.38), Cm(6.8)).table  
    shape.text = str(n) +'. ' + step
    shape.text_frame.paragraphs[0].font.size = Pt(20) 
    table.rows[0].height = Cm(1.54)
    table.columns[0].width  = Cm(1.4)
    table.columns[1].width = Cm(6.94)	 	 	 	 	 	 	 	 	 	 
    table.columns[2].width = Cm(3.58)
    table.columns[3].width = Cm(2.24)
    table.columns[4].width = Cm(2.45)
    table.columns[5].width = Cm(2.05)
    table.cell(0, 0).text = 'No'
    table.cell(0, 1).text = 'Protein Name'
    table.cell(0, 2).text = 'Concentration (mg/ml)'
    table.cell(0, 3).text = 'Volume (ml)'
    table.cell(0, 4).text = 'Amount (mg)'
    table.cell(0, 5).text = 'Yield (mg/L)'
    table.cell(0, 6).text = 'Buffer'
    table.cell(0, 7).text = 'Purity by SEC-HPLC (%)'
    table.cell(0, 8).text = 'Supernatant (mL)'
    if 'Protein A' in step or 'Ni' in step:
        table.columns[6].width = Cm(4.23)
        table.columns[7].width = Cm(3.98)
        table.columns[8].width = Cm(3.32)
        table.columns[9].width = Cm(1.07)
        table.columns[10].width = Cm(3.16)
        table.cell(0, 9).text = 'PI'
        table.cell(0, 10).text = 'Comments'
    else: 
        table.columns[6].width = Cm(1.89)
        table.columns[7].width = Cm(3.87)
        table.columns[8].width = Cm(3.23)
        table.columns[9].width = Cm(2.93)
        table.columns[10].width = Cm(1.07)
        table.columns[11].width = Cm(2.89)
        table.cell(0, 9).text = 'Recovery(%)'
        table.cell(0, 10).text = 'PI'
        table.cell(0, 11).text = 'Comments'
    for row,proteinName in enumerate(list(stepInfo.keys())[1:]):
        table.rows[row+1].height = Cm(1.05)
        table.cell(row+1,0).text = str(stepInfo[proteinName].get('No',''))
        table.cell(row+1,1).text = str(proteinName)
        for col in range(2,column):
            colName = table.cell(0, col).text.split('(')[0].strip()
            table.cell(row+1,col).text = \
            str(stepInfo[proteinName].get(colName,''))
    for cell in table.iter_cells():
        cell.margin_left = 0
        cell.margin_right = 0
        cell.margin_top = 0
        cell.margin_bottom = 0
        cell.vertical_anchor = 3
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.alignment = PP_ALIGN.CENTER            
    return prs

if __name__ == '__main__':
    filePath = os.path.abspath(os.path.dirname(__file__))
    template_path = os.path.join(filePath,'purificationTemplate.pptx')
    prs = Presentation(template_path)
    stepsDict = {
    1: {'PurificationStep': 'Protein A & Ni',
    'WXXX-hPro1.His - P1': { 
    'No': '1',
    'Concentration': '0.32',
    'Volume': '1.50',
    'Amount': '0.48',
    'Yield': '12.06',
    'Buffer': 'PBS',
    'Purity by SEC-HPLC': '99.62%',
    'Supernatant': '40',
    'PI': '6.32',
    'Recovery': 'NA',
    'MW': '147',
    'Comments': 'Three steps'},
    'WXXX-hPro1.His - P2' : {
    'No': '1',
    'Concentration': '1.68',
    'Volume': '2.30',
    'Amount': '3.87',
    'Yield': '96.80',
    'Buffer': 'PBS',
    'Purity by SEC-HPLC': '98.42%',
    'Supernatant': '40',
    'PI': '6.24',
    'Recovery': 'NA',
    'MW': '148',
    'Comments': 'Two steps'},
    'WXXX-hPro1.His - P3':{
    'No': '1',
    'Concentration': '0.90',
    'Volume': '1.50',
    'Amount': '1.35',
    'Yield': '33.70',
    'Buffer': 'PBS',
    'Purity by SEC-HPLC': '99.30%',
    'Supernatant': '40',
    'PI': '6.25',
    'Recovery': 'NA',
    'MW': '148',
    'Comments': 'Three steps'},
    'WXXX-BMK1': {
    'No': '2',
    'Concentration': '1.16',
    'Volume': '1.70',
    'Amount': '1.97',
    'Yield': '49.19',
    'Buffer': 'PBS',
    'Purity by SEC-HPLC': '96.30%',
    'Supernatant': '40',
    'PI': '6.31',
    'Recovery': 'NA',
    'MW': '147',
    'Comments': 'Two steps'},
    'WXXX-hPro1.His - P4': { 
    'No': '1',
    'Concentration': '0.32',
    'Volume': '1.50',
    'Amount': '0.48',
    'Yield': '12.06',
    'Buffer': 'PBS',
    'Purity by SEC-HPLC': '99.62%',
    'Supernatant': '40',
    'PI': '6.32',
    'Recovery': 'NA',
    'MW': '147',
    'Comments': 'Three steps'}
    },
    2: {'PurificationStep': 'Dialysis',
    'WXXX-hPro1.His - P5' : {
    'No': '1',
    'Concentration': '1.68',
    'Volume': '2.30',
    'Amount': '3.87',
    'Yield': '96.80',
    'Buffer': 'PBS',
    'Purity by SEC-HPLC': '98.42%',
    'Supernatant': '40',
    'PI': '6.24',
    'Recovery': 'NA',
    'MW': '148',
    'Comments': 'Two steps'},
    'WXXX-hPro1.His - P6':{
    'No': '1',
    'Concentration': '0.90',
    'Volume': '1.50',
    'Amount': '1.35',
    'Yield': '33.70',
    'Buffer': 'PBS',
    'Purity by SEC-HPLC': '99.30%',
    'Supernatant': '40',
    'PI': '6.25',
    'Recovery': 'NA',
    'MW': '148',
    'Comments': 'Three steps'},
    'WXXX-BMK2': {
    'No': '2',
    'Concentration': '1.16',
    'Volume': '1.70',
    'Amount': '1.97',
    'Yield': '49.19',
    'Buffer': 'PBS',
    'Purity by SEC-HPLC': '96.30%',
    'Supernatant': '40',
    'PI': '6.31',
    'Recovery': 'NA',
    'MW': '147',
    'Comments': 'Two steps'}
    },
    3: {'PurificationStep': 'SEC',
    'WXXX-hPro1.His - P7': { 
    'No': '1',
    'Concentration': '0.32',
    'Volume': '1.50',
    'Amount': '0.48',
    'Yield': '12.06',
    'Buffer': 'PBS',
    'Purity by SEC-HPLC': '99.62%',
    'Supernatant': '40',
    'PI': '6.32',
    'Recovery': 'NA',
    'MW': '147',
    'Comments': 'Three steps'},
    'WXXX-hPro1.His - P8' : {
    'No': '1',
    'Concentration': '1.68',
    'Volume': '2.30',
    'Amount': '3.87',
    'Yield': '96.80',
    'Buffer': 'PBS',
    'Purity by SEC-HPLC': '98.42%',
    'Supernatant': '40',
    'PI': '6.24',
    'Recovery': 'NA',
    'MW': '148',
    'Comments': 'Two steps'},
    'WXXX-hPro1.His - P9':{
    'No': '1',
    'Concentration': '0.90',
    'Volume': '1.50',
    'Amount': '1.35',
    'Yield': '33.70',
    'Buffer': 'PBS',
    'Purity by SEC-HPLC': '99.30%',
    'Supernatant': '40',
    'PI': '6.25',
    'Recovery': 'NA',
    'MW': '148',
    'Comments': 'Three steps'},
    'WXXX-BMK3': {
    'No': '2',
    'Concentration': '1.16',
    'Volume': '1.70',
    'Amount': '1.97',
    'Yield': '49.19',
    'Buffer': 'PBS',
    'Purity by SEC-HPLC': '96.30%',
    'Supernatant': '40',
    'PI': '6.31',
    'Recovery': 'NA',
    'MW': '147',
    'Comments': 'Two steps'}}}
    prs = purificationStepPageMake(prs, stepsDict)
    prs.save('purificationStepPageTest.pptx')