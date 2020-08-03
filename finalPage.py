# -*- coding: utf-8 -*-
"""
Created on Wed Jun 30 10:46:15 2020
@author: tao_jingfen
"""
from pptx import Presentation
import math
from pptx.util import Cm
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import os

def purificationFinalPageMake(prs: str, finalDict: dict) -> str:
    """
    Create the final protein purification summary page of purification auto-report PPT
    
    --Args:
        prs: presentation object loaded from *pptx* by template pptx
        finalDict: the final step protein info dict,
        e.g. { 
        'WXXX-hPro1.His - P1': 
        {
            'No':'1',
            'Concentration': 'XXX',
            'Volume': 'XXX',
            'Amount': 'XXX',
            'Yield': 'XXX',
            'Buffer': 'XXX',
            'Purity by SEC-HPLC': 'XXX',
            'Supernatant':'XXX',
            'MW':'XXX',
            'PI': 'XXX',
            'Comments': 'XXX'
        },
        'WXXX-hPro1.His - P2': 
        {
            'No':'1',
            'Concentration': 'XXX',
            'Volume': 'XXX',
            'Amount': 'XXX',
            'Yield': 'XXX',
            'Buffer': 'XXX',
            'Purity by SEC-HPLC': 'XXX',
            'Supernatant':'XXX',
            'MW':'XXX',
            'PI': 'XXX',
            'Comments': 'XXX'
        },
        'WXXX-BMK1': 
        {
            'No':'2',
            'Concentration': 'XXX',
            'Volume': 'XXX',
            'Amount': 'XXX',
            'Yield': 'XXX',
            'Buffer': 'XXX',
            'Purity by SEC-HPLC': 'XXX',
            'Supernatant':'XXX',
            'MW':'XXX',
            'PI': 'XXX',
            'Comments': 'XXX'
        },
         ...}
    --Returns:
        prs：presentation object could be used to save pptx
        10 samples each page at most
        more samples will generate another page
    """
    sampleNumber = len(finalDict.keys())
    #for page in range(1,max(math.ceil(sampleNumber/10)+1,2)):
    for page in range(1,math.ceil(sampleNumber/10)+1):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        shapes = slide.shapes
        shapes.title.text = 'Characterization'
        subtitle = shapes.title.text_frame.add_paragraph()
        subtitle.text = 'Final '
        subtitle.font.italic = True
        subtitle1 = subtitle.add_run()
        subtitle1.text = '- Protein Purification Summary'
        subtitle1.font.italic = True
        subtitle1.font.bold = False
        if page * 10 < sampleNumber:
            rowNumber = 10
        else:
            rowNumber = sampleNumber-(page-1)*10
        table = shapes[1].insert_table(rows=rowNumber+1, cols=12).table
        table.rows[0].height = Cm(1.92)
        table.columns[0].width = Cm(1.23)
        table.columns[1].width = Cm(6.18)
        table.columns[2].width = Cm(3.58)
        table.columns[3].width = Cm(2.17)
        table.columns[4].width = Cm(2.37)
        table.columns[5].width = Cm(1.99)
        table.columns[6].width = Cm(1.88)
        table.columns[7].width = Cm(3.85)
        table.columns[8].width = Cm(3.21)
        table.columns[9].width = Cm(1.78)
        table.columns[10].width = Cm(1.47)
        table.columns[11].width = Cm(3.89)
        table.cell(0, 0).text = 'No'
        table.cell(0, 1).text = 'Protein Name'
        table.cell(0, 2).text = 'Concentration (mg/ml)'
        table.cell(0, 3).text = 'Volume (ml)'
        table.cell(0, 4).text = 'Amount (mg)'
        table.cell(0, 5).text = 'Yield (mg/L)'
        table.cell(0, 6).text = 'Buffer'
        table.cell(0, 7).text = 'Purity by SEC-HPLC (%)'
        table.cell(0, 8).text = 'Supernatant (mL)'
        table.cell(0, 9).text = 'MW (kDa)'
        table.cell(0, 10).text = 'PI'
        table.cell(0, 11).text = 'Comments'
        for row in range(0,rowNumber):
            table.rows[row+1].height = Cm(1.2)
            proteinName = list(finalDict.keys())[(page-1)*10+row]
            table.cell(row+1,0).text = str(finalDict[proteinName].get('No',''))
            table.cell(row+1,1).text = str(proteinName)
            for col in range(2,12):
                colName = table.cell(0, col).text.split('(')[0].strip()
                table.cell(row+1,col).text = str(finalDict[proteinName].get(colName,''))
        for cell in table.iter_cells():
            cell.margin_left = 0
            cell.margin_right = 0
            cell.margin_top = 0
            cell.margin_bottom = 0
            cell.vertical_anchor = 3 ## 垂直居中
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.alignment = PP_ALIGN.CENTER ## 2  水平居中 
    return prs

if __name__ == '__main__':
    filePath = os.path.abspath(os.path.dirname(__file__))
    template_path = os.path.join(filePath,'purificationTemplate.pptx')
    prs = Presentation(template_path)
    finalDict = {
    'WXXX-hPro1.His - P1': { 
    'No': '1',
    'Concentration': '0.32',
    'Volume': '1.50',
    'Amount': '0.48',
    'Yield': '12.06',
    'Buffer': 'PBS',
    'Purity by SEC-HPLC': '99.62%',
    'Supernatant': '40',
    'PI': '6.32',
    'Recovery': 'NA',
    'MW': '147',
    'Comments': 'Three steps'},
    'WXXX-hPro1.His - P2' : {
    'No': '1',
    'Concentration': '1.68',
    'Volume': '2.30',
    'Amount': '3.87',
    'Yield': '96.80',
    'Buffer': 'PBS',
    'Purity by SEC-HPLC': '98.42%',
    'Supernatant': '40',
    'PI': '6.24',
    'Recovery': 'NA',
    'MW': '148',
    'Comments': 'Two steps'},
    'WXXX-hPro1.His - P3':{
    'No': '1',
    'Concentration': '0.90',
    'Volume': '1.50',
    'Amount': '1.35',
    'Yield': '33.70',
    'Buffer': 'PBS',
    'Purity by SEC-HPLC': '99.30%',
    'Supernatant': '40',
    'PI': '6.25',
    'Recovery': 'NA',
    'MW': '148',
    'Comments': 'Three steps'},
    'WXXX-BMK1': {
    'No': '2',
    'Concentration': '1.16',
    'Volume': '1.70',
    'Amount': '1.97',
    'Yield': '49.19',
    'Buffer': 'PBS',
    'Purity by SEC-HPLC': '96.30%',
    'Supernatant': '40',
    'PI': '6.31',
    'Recovery': 'NA',
    'MW': '147',
    'Comments': 'Two steps'},
    'WXXX-hPro1.His - P4': { 
    'No': '1',
    'Concentration': '0.32',
    'Volume': '1.50',
    'Amount': '0.48',
    'Yield': '12.06',
    'Buffer': 'PBS',
    'Purity by SEC-HPLC': '99.62%',
    'Supernatant': '40',
    'PI': '6.32',
    'Recovery': 'NA',
    'MW': '147',
    'Comments': 'Three steps'},
    'WXXX-hPro1.His - P5' : {
    'No': '1',
    'Concentration': '1.68',
    'Volume': '2.30',
    'Amount': '3.87',
    'Yield': '96.80',
    'Buffer': 'PBS',
    'Purity by SEC-HPLC': '98.42%',
    'Supernatant': '40',
    'PI': '6.24',
    'Recovery': 'NA',
    'MW': '148',
    'Comments': 'Two steps'},
    'WXXX-hPro1.His - P6':{
    'No': '1',
    'Concentration': '0.90',
    'Volume': '1.50',
    'Amount': '1.35',
    'Yield': '33.70',
    'Buffer': 'PBS',
    'Purity by SEC-HPLC': '99.30%',
    'Supernatant': '40',
    'PI': '6.25',
    'Recovery': 'NA',
    'MW': '148',
    'Comments': 'Three steps'},
    'WXXX-BMK2': {
    'No': '2',
    'Concentration': '1.16',
    'Volume': '1.70',
    'Amount': '1.97',
    'Yield': '49.19',
    'Buffer': 'PBS',
    'Purity by SEC-HPLC': '96.30%',
    'Supernatant': '40',
    'PI': '6.31',
    'Recovery': 'NA',
    'MW': '147',
    'Comments': 'Two steps'},
    'WXXX-hPro1.His - P7': { 
    'No': '1',
    'Concentration': '0.32',
    'Volume': '1.50',
    'Amount': '0.48',
    'Yield': '12.06',
    'Buffer': 'PBS',
    'Purity by SEC-HPLC': '99.62%',
    'Supernatant': '40',
    'PI': '6.32',
    'Recovery': 'NA',
    'MW': '147',
    'Comments': 'Three steps'},
    'WXXX-hPro1.His - P8' : {
    'No': '1',
    'Concentration': '1.68',
    'Volume': '2.30',
    'Amount': '3.87',
    'Yield': '96.80',
    'Buffer': 'PBS',
    'Purity by SEC-HPLC': '98.42%',
    'Supernatant': '40',
    'PI': '6.24',
    'Recovery': 'NA',
    'MW': '148',
    'Comments': 'Two steps'},
    'WXXX-hPro1.His - P9':{
    'No': '1',
    'Concentration': '0.90',
    'Volume': '1.50',
    'Amount': '1.35',
    'Yield': '33.70',
    'Buffer': 'PBS',
    'Purity by SEC-HPLC': '99.30%',
    'Supernatant': '40',
    'PI': '6.25',
    'Recovery': 'NA',
    'MW': '148',
    'Comments': 'Three steps'},
    'WXXX-BMK3': {
    'No': '2',
    'Concentration': '1.16',
    'Volume': '1.70',
    'Amount': '1.97',
    'Yield': '49.19',
    'Buffer': 'PBS',
    'Purity by SEC-HPLC': '96.30%',
    'Supernatant': '40',
    'PI': '6.31',
    'Recovery': 'NA',
    'MW': '147',
    'Comments': 'Two steps'}}
    prs = purificationFinalPageMake(prs, finalDict)
    prs.save('purificationFinalPageTest.pptx')
