'''
@Author: tao_jingfen
@Date: 2019-07-05 11:17:24
@LastEditors: tao_jingfen
@LastEditTime: 2020-07-22 14:22:34
'''

import os
import tkinter
from tkinter import messagebox
from tkinter import ttk
import tkinter.filedialog
from pptx import Presentation

from readExcel import excel2Dict
from coverPage import purificationCoverPageMake
from finalPage import purificationFinalPageMake
from processPage import purificationProcessPageMake
from stepPage import purificationStepPageMake
from sdsPage import purificationSdsPageMake
from hplcPage import purificationHplcPageMake

import warnings
warnings.filterwarnings("ignore")

class AutoReport(object):
    def __init__(self):
        self.tk = tkinter.Tk()
        self.tk.geometry("400x200")
        self.tk.resizable(0, 0)
        self.filePath = os.path.abspath(os.path.dirname(__file__))
        self.tadPic = os.path.join(self.filePath,'bid.ico')
        self.tk.iconbitmap(self.tadPic)
        self.tk.title("BID AutoReport System")

        self.report_type = tkinter.StringVar() 
        self.type_label = tkinter.Label(self.tk,  text="Report Type:", font=('Calibri', '12', 'bold'))
        self.type_combobox = ttk.Combobox(self.tk, textvariable=self.report_type, state='readonly')
        self.type_combobox['values'] = ('Purification Report', 'Expression Report', 'Analysis Report')
        self.type_combobox.current(0)
        
        self.path_label = tkinter.Label(self.tk,  text="Excel Path:", font=('Calibri', '12', 'bold'))
        self.path_text = tkinter.Text(self.tk,height=4,width=36,font=('Calibri','11'))
        self.path_button = ttk.Button(self.tk, text="...",command=self._select_path,width=3)

        self.start_button = tkinter.Button(self.tk,  text="Generate Report", font=('Calibri', '12', 'bold'), command=self.generate_report)
        self.copyright_label = tkinter.Label(self.tk, text='Copyright ©2020. Powered by BID BioIT.', font=('Times New Roman', '10'))
        self.filePath = os.getcwd()
        
    def item_pack(self):
        self.type_label.grid(row = 0, ipadx=2, ipady=3, pady=5,sticky = tkinter.W)
        self.type_combobox.grid(row = 0, column = 1, columnspan=2,padx=2,ipady=2,sticky = tkinter.W + tkinter.E)
        self.path_label.grid(row = 1,  ipadx=2, padx=3, pady=5, sticky = tkinter.W + tkinter.E)
        self.path_text.grid(row = 1, column = 1, padx=2, ipady=3, pady=5, sticky = tkinter.W + tkinter.E)
        self.path_button.grid(row = 1, column = 2,padx=3,sticky = tkinter.W + tkinter.E)
        self.start_button.grid(row = 2,column = 1, pady=3)
        self.copyright_label.grid(row = 3,column = 0,columnspan=3,pady=3,sticky = tkinter.W + tkinter.E)

    @staticmethod
    def mess(mes: str):
        """
        output messagebox.showwarning with the mes
        """
        messagebox.showwarning('Warnings:', mes)
        
    def generate_report(self):
        try:
            self.workdir = os.path.dirname(self.excelPath)
            if self.report_type.get() == 'Purification Report':
                self.purification_report()
                tkinter.messagebox.showinfo('Success','Generate Report Finished!')
        except Exception as mes:
            self.mess(mes)  

    def purification_report(self):
        self.template_path = os.path.join(self.filePath,'purificationTemplate.pptx')
        prs = Presentation(self.template_path)
        os.chdir(self.workdir)
        (projectName,
         date),finalDict,processList,stepsDict,(supernatant,
             sdsList,sdsTableDict),hplcList,proDict = excel2Dict(self.excelPath)
        prs = purificationCoverPageMake(prs, projectName, date)
        prs = purificationFinalPageMake(prs, finalDict)
        prs = purificationProcessPageMake(prs, processList)
        prs = purificationStepPageMake(prs, stepsDict)
        prs = purificationSdsPageMake(prs, supernatant,sdsList,sdsTableDict)
        prs = purificationHplcPageMake(prs, hplcList, proDict)
        numPlaceholder = prs.slide_layouts[1].placeholders[0]
        for num in range(1,len(prs.slides)):
            slide = prs.slides[num]
            slide.shapes.clone_placeholder(numPlaceholder)
            slide.shapes[-1].text = str(num + 1)
        prs.save(''.join([date.split('/')[-1]]+date.split('/')[0:2]) + ' ' + 
                 projectName + ' Purification report.pptx')
        #prs.save(projectName + '_' + ''.join([date.split('/')[-1]]+date.split('/')[0:2]) + '.pptx')

    def _select_path(self):
        self.excelPath = tkinter.filedialog.askopenfilename()
        self.excelPath = self.excelPath.replace("/","\\\\")
        self.path_text.delete(0.0, tkinter.END)
        self.path_text.insert('0.0', self.excelPath)

if __name__ == "__main__":
    tmp = AutoReport()
    tmp.item_pack()
    tmp.tk.mainloop()