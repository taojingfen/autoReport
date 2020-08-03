'''
@Author: tao_jingfen
@Date: 2020-06-30 13:41:44
@LastEditors: tao_jingfen
@LastEditTime: 2020-07-03 10:20:31
'''

from pptx import Presentation
import math
from pptx.util import Cm
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

def purificationProcessPageMake(prs: str, processList: list) -> str:
    """
    Create the prorein purification process page of purification auto-report PPT
    
    --Args:
        prs: presentation object loaded from *pptx* by template pptx
        processList: every protein process info list,
        e.g. [
        ['Protein # 1, 3', 'Protein A', 'Dialysis', 'SEC','CEX','Dialysis','Filtration & Storage'],
        ['Protein # 2, 4', 'Protein A', 'Dialysis', 'SEC','Filtration & Storage']   
        ]
    --Returns:
        prs：presentation object could be used to save pptx
        3 process types each page at most
        more will generate another page
    """
    processNumber = len(processList)
    #for page in range(1,max(math.ceil(processNumber/3)+1,2)):
    for page in range(1,math.ceil(processNumber/3)+1):
        slide = prs.slides.add_slide(prs.slide_layouts[2])
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = 'Characterization'
        subtitle = title_shape.text_frame.add_paragraph()
        subtitle.text = 'Protein Purification Process'
        subtitle.font.italic = True
        subtitle.font.bold = False
        if page * 3 < processNumber:
            rowNumber = 3
        else:
            rowNumber = processNumber-(page-1)*3
        left, width, height = Cm(2.3), Cm(31.5), Cm(4.28)
        for row in range(0,rowNumber):
            top = Cm([3,8.5,14][row])
            shape = shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
            shape.shadow.inherit = False
            shape.fill.solid()  ##实色填充
            shape.fill.fore_color.rgb = RGBColor(208,216,232)
            shape.line.fill.background()  ##透明边缘
            typeNumber = len(processList[(page-1)*3+row])
            margin = round((27.9-typeNumber*3.6)/(typeNumber-1),2)
            small_left, small_top,small_width, small_height = Cm(3), top+Cm(1.05), Cm(3.6), Cm(2.1)
            for n in range(typeNumber):
                shape = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, small_left, small_top, small_width, small_height)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(79,129,189)
                shape.line.color.rgb = RGBColor(255,255,255)
                shape.line.width = Cm(0.07)
                shape.shadow.inherit = False
                shape.text = processList[(page-1)*3+row][n]
                shape.text_frame.paragraphs[0].font.size = Pt(14)
                shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                small_left = small_left + small_width + Cm(margin)
    return prs

if __name__ == '__main__':
    filePath = os.path.abspath(os.path.dirname(__file__))
    template_path = os.path.join(filePath,'purificationTemplate.pptx')
    prs = Presentation(template_path)
    processList = [
        ['Protein # 1, 3', 'Protein A', 'Dialysis', 'SEC','CEX','Dialysis','Filtration & Storage'],
        ['Protein # 2, 4', 'Protein A', 'Dialysis', 'SEC','Filtration & Storage'],
        ['Protein # 5', 'Ni', 'Dialysis', 'SEC','CEX','Dialysis','Filtration & Storage'],
        ['Protein # 6', 'Ni', 'Dialysis', 'SEC','Filtration & Storage']   
        ]
    prs = purificationProcessPageMake(prs, processList)
    prs.save('purificationProcessPageTest.pptx')