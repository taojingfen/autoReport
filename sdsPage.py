# -*- coding: utf-8 -*-
'''
@Author: tao_jingfen
@Date: 2020-07-01 14:34:07
@LastEditors: tao_jingfen
@LastEditTime: 2020-07-21 11:34:43
'''
from pptx import Presentation
import os
import math
from pptx.util import Cm
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR

filePath = os.path.abspath(os.path.dirname(__file__))

def purificationSdsPageMake(prs: str, supernatant: str,sdsList: list, sdsTableDict: dict) -> str:
    """
    Create protein SDS-PAGE page of purification auto-report PPT
    
    --Args:
        prs: presentation object loaded from *pptx* by template pptx
        supernatant: the supernatant volume(ml), e.g. "40"
        sdsList: the list of every SDS-PAGE step info,
        [stepName, SDSpictureName, SDSTableKeyName,SDSelnInfo,SDSConclusionInfo]
        e.g. [
            ['Protein A', 'SupernatantDialysis_SDS-PAGE.png', 'SupernatantDialysis','SDS_ELN','SDS_Conclusion'],
            ['Dialysis', 'SEC_SDS-PAGE.png', 'SEC','SDS_ELN','SDS_Conclusion']
            ...
            ]
        sdsTableDict:  the dict of SDS table info,
        e.g.
        {
        'SupernatantDialysis': {
            'Lane': ['S1', 'S2','1','2'],
            'Protein name': ['WXXX-T6U9.E17-1.uIgG1 (supernatant)',
               'WXXX-U9T6.E17-1.uIgG1 (supernatant)',
               'WXXX-T6U9.E17R-1.uIgG1',
               'WXXX-U9T6.E17R-1.uIgG1'],
            'MW(kDa)': ['147',
               '148',
               '22+53+50+23',
               '22+53+49+23']},
        'SEC': {
            'Lane': ['1', '2', '3'],
            'Protein name': ['WXXX-T6U9.E17-1.uIgG1',
               'WXXX-U9T6.E17-1.uIgG1',
               'WXXX-T6U9.E17R-1.uIgG1'],
            'MW(kDa)': ['147',
               '148',
               '148']},
        'CEX': {
            'Lane': ['1', '2'],
            'Protein name': ['WXXX-T6U9.E17-1.uIgG1','WXXX-T6U9.E17R-1.uIgG1],
            'MW(kDa)': ['147', '148', '22+53+49+23', '22+53+50+23']}}
    --Returns:
        prs：presentation object could be used to save pptx
        8 samples in SDS table each page at most
        more samples will generate another page
    """
    textList = [supernatant + 'mL Protein Supernatant']
    for stepSds in sdsList:
        if stepSds[0] == 'Supernatant':
            vol = '13 μL on SDS-PAGE'
        else:
            textList.append(stepSds[0])
            vol = '2 μg on SDS-PAGE'
        if stepSds[1] == '-' or stepSds[1] == '':
            continue
        prs = sdsMake(prs,stepSds[0],textList + [vol,'Staining'],stepSds[1],sdsTableDict.get(stepSds[2],{}),stepSds[3],stepSds[4])
    return prs

def sdsMake(prs,stepTitle,textList,SDS_Picture,SDS_TableDict,SDS_ELN,SDS_Conclusion):
    """
    Create single protein SDS-PAGE page
    
    --Args:
        prs: presentation object loaded from *pptx* by template pptx
        stepTitle: the step that show in the title, e.g. "Supernatant & Dialysis"
        textList: the textList use to generate step arrow,
        e.g. ['40mL Protein Supernatant', 'Protein A', 'Dialysis', '13 μL / 2 μg on SDS-PAGE', 'Staining']
        SDS_Picture: the SDS-PAGE picture of this step, e.g. "SupernatantDialysis_SDS-PAGE.png"
        SDS_TableDict:    
        e.g. {
            'Lane': ['S1', 'S2','1','2'],
            'Protein name': ['WXXX-T6U9.E17-1.uIgG1 (supernatant)',
               'WXXX-U9T6.E17-1.uIgG1 (supernatant)',
               'WXXX-T6U9.E17R-1.uIgG1',
               'WXXX-U9T6.E17R-1.uIgG1'],
            'MW(kDa)': ['147',
               '148',
               '22+53+50+23',
               '22+53+49+23']}
        SDS_ELN: the ELN info of SDS-PAGE setp, e.g. "WXBIOXXX/XXXX/20200103"
        SDS_Conclusion: the conclusion of SDS-PAGE that split by '|',
        e.g. Expected reducing and non-reducing bands were visible from the gel.|Minor bands were observed in Lane # 1, 2, 3, 4. 
    --Returns:
        prs：presentation object could be used to save pptx
        12 samples in SDS_table each page at most
        more samples will generate another page
    """
    tableRowNum = len(SDS_TableDict.get('Lane',[]))
    for page in range(1,max(math.ceil(tableRowNum/12+1),2)):
        slide = prs.slides.add_slide(prs.slide_layouts[3])
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = 'Characterization'
        subtitle = title_shape.text_frame.add_paragraph()
        subtitle.text = 'SDS-PAGE Results – ' + stepTitle
        subtitle.font.italic = True
        subtitle.font.bold = False
        left, top, width, height = Cm(2.3), Cm(3),Cm(31.5), Cm(4.28)
        shape = shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
        shape.shadow.inherit = False
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(208,216,232)
        shape.line.fill.background()
        stepNumber = len(textList)
        margin = round((27.9-stepNumber*3.6)/(stepNumber-1),2)
        small_left, small_top,small_width, small_height = Cm(3), top+Cm(1.05), Cm(3.6), Cm(2.1)
        for text in textList:
            shape = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, small_left, small_top, small_width, small_height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(79,129,189)
            shape.line.color.rgb = RGBColor(255,255,255)
            shape.line.width = Cm(0.07)
            shape.shadow.inherit = False
            shape.text = str(text)
            shape.text_frame.paragraphs[0].font.size = Pt(14)
            shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            small_left = small_left + small_width + Cm(margin)
        if page * 12 <= tableRowNum:
                rows = 12
        else:
            rows = tableRowNum-(page-1)*12
        laneList = SDS_TableDict['Lane'][(page-1)*12:(page-1)*12+rows] + ['M']
        sdsPicture(shapes,SDS_Picture,laneList)
        cols = len(SDS_TableDict.keys())
        table = shapes.add_table(rows + 2, cols, Cm(16.6), Cm(6.6), Cm(14.2), Cm(1)).table
        table.rows[0].height = Cm(1.13)
        table.columns[0].width  = Cm(1.54)
        table.columns[1].width = Cm(9.5)	 	 	 	 	 	 	 	 	 	 
        table.columns[2].width = Cm(3.8)
        table.cell(0, 0).text = 'Lane'
        table.cell(0, 1).text = 'Protein name'
        table.cell(0, 2).text = 'MW(kDa)'
        for col in range(cols):
            for row in range(1,rows+1):
                table.rows[row].height = Cm(0.71)
                table.cell(row,col).text = SDS_TableDict[table.cell(0, col).text][(page-1)*12+row-1]
        table.rows[rows+1].height = Cm(0.71)
        table.cell(rows+1,0).text = 'M'
        table.cell(rows+1, 1).merge(table.cell(rows+1, cols-1))
        table.cell(rows+1, 1).text = 'PageRuler™ Unstained Protein  Ladder'
        for cell in table.iter_cells():
            cell.margin_left = 0
            cell.margin_right = 0
            cell.margin_top = 0
            cell.margin_bottom = 0
            cell.vertical_anchor = 3 
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.alignment = PP_ALIGN.CENTER       
        shape = shapes[0]
        shape.left,shape.top,shape.width,shape.height = Cm(2),Cm(16.2),Cm(31.5),Cm(3) 
        shape.text = 'Conclusion'
        shape.text_frame.paragraphs[0].font.size = Pt(18)
        shape.text_frame.paragraphs[0].font.bold = True
        for text in SDS_Conclusion.split('|'):
            new_paragraph = shape.text_frame.add_paragraph()
            new_paragraph.text = text
            new_paragraph.font.bold = False
            new_paragraph.font.size = Pt(14)
            new_paragraph.level = 1
        shape = shapes.add_textbox(Cm(28.5), Cm(18.8), Cm(3), Cm(1.11))
        shape.text = SDS_ELN
        shape.text_frame.paragraphs[0].font.size = Pt(10)
    return prs

def sdsPicture(shapes,picture,laneList):
    shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(2.5), Cm(6.6), Cm(13.6), Cm(8.8))
    shape.shadow.inherit = False
    shape.fill.background()
    shape.line.color.rgb = RGBColor(56,93,138)
    shape.line.width = Cm(0.07)
    markerPic = os.path.join(filePath,'marker.jpg')
    shapes.add_picture(markerPic,Cm(2.7), Cm(7.8))
    shape = shapes.add_textbox(Cm(5.7), Cm(6.6), Cm(3.7), Cm(0.9))
    shape.text = 'Non-reducing'
    shape.text_frame.paragraphs[0].font.size = Pt(14)
    shape = shapes.add_textbox(Cm(10.1), Cm(6.6), Cm(2.4), Cm(0.9))
    shape.text = 'Reducing'
    shape.text_frame.paragraphs[0].font.size = Pt(14)
    shape = shapes.add_textbox(Cm(3.9), Cm(7.4), Cm(8.8), Cm(0.9))
    shape.text = '    '.join(laneList)
    shape.text_frame.paragraphs[0].font.size = Pt(14)
    connector = shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Cm(4.2),Cm(7.5),Cm(9.8),Cm(7.5))
    connector.line.color.rgb = RGBColor(0,0,0)
    connector.shadow.inherit = False
    connector = shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Cm(10),Cm(7.5),Cm(12.4),Cm(7.5))
    connector.line.color.rgb = RGBColor(0,0,0)
    connector.shadow.inherit = False
    shapes.add_picture(picture,Cm(4.18),Cm(8.1),Cm(9),Cm(5.82))
    shape = shapes.add_textbox(Cm(3.3), Cm(14.5), Cm(10.3), Cm(0.9))
    shape.text = 'Gel info: NuPAGE, Novex 4-12% Bis-Tris Gel'
    shape.text_frame.paragraphs[0].font.size = Pt(16)
    arrowPic = os.path.join(filePath,'arrow.png')
    for top in [Cm(8.56),Cm(10.1),Cm(11.3)]:
        shapes.add_picture(arrowPic,Cm(12.95),top)
    shape = shapes.add_textbox(Cm(14), Cm(8.4), Cm(1.7), Cm(0.9))
    shape.text = '150 kDa'
    shape.text_frame.paragraphs[0].font.size = Pt(14)
    shape = shapes.add_textbox(Cm(14), Cm(10), Cm(1.7), Cm(0.9))
    shape.text = '50 kDa'
    shape.text_frame.paragraphs[0].font.size = Pt(14)
    shape = shapes.add_textbox(Cm(14), Cm(11.2), Cm(1.7), Cm(0.9))
    shape.text = '25 kDa'
    shape.text_frame.paragraphs[0].font.size = Pt(14)

if __name__ == '__main__':
    template_path = os.path.join(filePath,'purificationTemplate.pptx')
    prs = Presentation(template_path)
    supernatant = '40'
    sdsList = [
        ['Protein A','-','-','-','-'],
        ['Dialysis','SupernatantDialysis_SDS-PAGE.png','SupernatantDialysis','WXBIOXXX/XXXX/20200103','Expected reducing and non-reducing bands were visible from the gel.|Minor bands were observed in Lane # 1, 2, 3, 4.'], 
        ['SEC','SEC_SDS-PAGE.png','SEC','WXBIOXXX/XXXX/20200106','Expected reducing and non-reducing bands were visible from the gel.|Minor bands were observed in Lane # 1, 3, 4. '],
        ['CEX','CEX_SDS-PAGE.png','CEX','WXBIOXXX/XXXX/20200110','Expected reducing and non-reducing bands were visible from the gel.|Minor bands were observed in Lane # 1, 2. '],
        ['Dialysis','-','-','-','-']
    ]
    sdsTableDict = {
        'SupernatantDialysis': {
            'Lane': ['S1','S2','S3','S4','1','2','3','4','5','6','7','8'],
            'Protein name': ['WXXX-T6U9.E17-1.uIgG1 (supernatant)',
                'WXXX-U9T6.E17-1.uIgG1 (supernatant)',
                'WXXX-T6U9.E17R-1.uIgG1 (supernatant)',
                'WXXX-U9T6.E17R-1.uIgG1 (supernatant)',
                'WXXX-T6U9.E17-1.uIgG1',
                'WXXX-U9T6.E17-1.uIgG1',
                'WXXX-T6U9.E17R-1.uIgG1',
                'WXXX-U9T6.E17R-1.uIgG1',
                'WXXX-T6U9.E17-1.uIgG1',
                'WXXX-U9T6.E17-1.uIgG1',
                'WXXX-T6U9.E17R-1.uIgG1',
                'WXXX-U9T6.E17R-1.uIgG1'],
            'MW(kDa)': ['147','148','148','147','147','148','148','147',
                '22+53+49+23',
                '22+53+49+23',
                '22+53+50+23',
                '22+53+49+23']},
        'SEC': {
            'Lane': ['1', '2', '3', '4', '5', '6', '7', '8'],
            'Protein name': ['WXXX-T6U9.E17-1.uIgG1',
                'WXXX-U9T6.E17-1.uIgG1',
                'WXXX-T6U9.E17R-1.uIgG1',
                'WXXX-U9T6.E17R-1.uIgG1',
                'WXXX-T6U9.E17-1.uIgG1',
                'WXXX-U9T6.E17-1.uIgG1',
                'WXXX-T6U9.E17R-1.uIgG1',
                'WXXX-U9T6.E17R-1.uIgG1'],
            'MW(kDa)': ['147','148','148','147',
                '22+53+49+23',
                '22+53+49+23',
                '22+53+50+23',
                '22+53+49+23']},
        'CEX': {
            'Lane': ['1', '2', '3', '4'],
            'Protein name': ['WXXX-T6U9.E17-1.uIgG1',
                'WXXX-T6U9.E17R-1.uIgG1',
                'WXXX-T6U9.E17-1.uIgG1',
                'WXXX-T6U9.E17R-1.uIgG1'],
            'MW(kDa)': ['147', '148', '22+53+49+23', '22+53+50+23']}
    }
    prs = purificationSdsPageMake(prs,supernatant,sdsList,sdsTableDict)
    prs.save('purificationSdsPageTest.pptx')