# -*- coding: utf-8 -*-
'''
@Author: tao_jingfen
@Date: 2020-07-01 14:54:44
@LastEditors: tao_jingfen
@LastEditTime: 2020-07-30 10:47:30
'''

import os
import docx
import math
import fitz
import re
from pptx import Presentation
from pptx.util import Cm
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pdfminer.pdfparser import PDFParser, PDFDocument
from pdfminer.converter import PDFPageAggregator
from pdfminer.layout import LAParams, LTTextBoxHorizontal
from pdfminer.pdfinterp import PDFTextExtractionNotAllowed, PDFResourceManager, PDFPageInterpreter


def purificationHplcPageMake(prs: str, hplcList: list, proDict: dict) -> str:
    """
    Create protein SEC-HPLC page of purification auto-report PPT

    --Args:
        prs: presentation object loaded from *pptx* by template pptx
        supernatant: the supernatant volume(ml), e.g. "40"
        hplcList: the list of every SEC-HPLC step info,
        [stepName, HPLCfileName, HPLCelnInfo,HPLCconclusionInfo]
        e.g. [
            ['1','Protein A', 'Protein A.docx','HPLC_ELN','HPLC_Conclusion'],
            ['2','Dialysis', 'SEC.pdf', HPLC_ELN','HPLC_Conclusion']
            ...
            ]
        proDict: the protein No and protein purity Dict,
        {stepName:{proteinName:[proteinNo,purity],...}}
        e.g.{
            '1_Protein_A':{
                'WXXX-T6U9.E17-1.uIgG1': [1,'59.77%'],
                'WXXX-U9T6.E17-1.uIgG1': [2,'86.79%'],
                'WXXX-T6U9.E17R-1.uIgG1': [3,'85.68%'],
                'WXXX-U9T6.E17R-1.uIgG1': [4,'75.10%']
            },
            '2_Dialysis':{
                'WXXX-T6U9.E17-1.uIgG1': [1,'88.57%'],
                'WXXX-U9T6.E17-1.uIgG1': [2,'86.34%'],
                'WXXX-T6U9.E17R-1.uIgG1': [3,'87.80%'],
                'WXXX-U9T6.E17R-1.uIgG1': [4,'85.81%']
            },
            ...}
    --Returns:
        prs：presentation object could be used to save pptx
        4 pictures in HPLC each page at most
        more samples will generate another page
    """
    for stepHplc in hplcList:
        stepName = '_'.join((stepHplc[0],stepHplc[1].replace(' ','_')))
        stepProDict = proDict.get(stepName,{})
        prs = hplcMake(prs,stepProDict,*stepHplc[1:])
    return prs

def hplcMake(prs,stepProDict,step,HPLC_File,HPLC_ELN,HPLC_Conclusion):
    """
    Create single protein SEC-HPLC page
    
    --Args:
        prs: presentation object loaded from *pptx* by template pptx
        step: the step name that also show in the title, e.g. "Protein A "
        stepProDict: dict, the protein No and protein purity of the corresponding protein name
        e.g. {
            'WXXX-T6U9.E17-1.uIgG1': [1,'59.77%'],
            'WXXX-U9T6.E17-1.uIgG1': [2,'86.79%'],
            'WXXX-T6U9.E17R-1.uIgG1': [3,'85.68%'],
            'WXXX-U9T6.E17R-1.uIgG1': [4,'75.10%']
            }
        HPLC_File: the word file or pdf file that contain HPLC picture of this step, e.g. "SupernatantDialysis_SEC_HPLC.docx"
        HPLC_ELN: the ELN info of SEC-HPLC setp, e.g. "WXBIOXXX/XXXX/20200103"
        HPLC_Conclusion: the conclusion of SEC-HPLC that split by '|',
        e.g. Expected reducing and non-reducing bands were visible from the gel.|Minor bands were observed in Lane # 1, 2, 3, 4. 
    --Returns:
        prs：presentation object could be used to save pptx
        4 pictures in HPLC each page at most
        more samples will generate another page
    """
    if HPLC_File.strip() != '':
        if HPLC_File.split('.')[-1] == 'docx':
            picList = word2pic(HPLC_File)
        elif HPLC_File.split('.')[-1] == 'pdf':
            picList = pdf2pic(HPLC_File)
        else:
            raise ValueError("HPLC文件必须是docx或PDF")
        sampleNumber = len(picList)
        for page in range(1,max(math.ceil(sampleNumber/10)+1,2)):
            slide = prs.slides.add_slide(prs.slide_layouts[3])
            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.text = 'Characterization'
            subtitle = title_shape.text_frame.add_paragraph()
            subtitle.text = 'SEC-HPLC Results – '+ step
            subtitle.font.italic = True
            subtitle.font.bold = False
            if page * 4 < sampleNumber:
                rowNumber = 4
            else:
                rowNumber = sampleNumber-(page-1)*4
            leftList = [Cm(5.2),Cm(18.5)]
            topList = [Cm(3.2), Cm(8.8)] 
            for n in range(rowNumber):
                sampleIndex = (page-1)*4+n
                sampleName = picList[sampleIndex].rsplit(".", 1)[0]
                sampleNo = stepProDict.get(sampleName,['1','100%'])[0]
                purity = stepProDict.get(sampleName,['1','100%'])[1]
                top = topList[math.ceil((n+1)/2)-1]
                left = leftList[n%2]
                shape = shapes.add_textbox(left, top, Cm(12), Cm(0.6))
                shape.text = str(sampleNo) + '.' + sampleName + \
                        ', ' + purity
                shape.text_frame.paragraphs[0].font.size = Pt(14)
                shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                shape = shapes.add_picture(picList[sampleIndex],left,top + Cm(0.8),Cm(12),Cm(4))
                shape.line.color.rgb = RGBColor(79,129,189)
                os.remove(picList[sampleIndex])
            shape = shapes[0]
            shape.left,shape.top,shape.width,shape.height = Cm(2.3),Cm(14.2),Cm(31.5),Cm(3) 
            shape.text = 'Conclusion'
            shape.text_frame.paragraphs[0].font.size = Pt(18)
            shape.text_frame.paragraphs[0].font.bold = True
            for text in HPLC_Conclusion.split('|'):
                new_paragraph = shape.text_frame.add_paragraph()
                new_paragraph.text = text
                new_paragraph.font.bold = False
                new_paragraph.font.size = Pt(14)
                new_paragraph.level = 1
            shape = shapes.add_textbox(Cm(28.5), Cm(18.8), Cm(3), Cm(1.11))
            shape.text = HPLC_ELN
            shape.text_frame.paragraphs[0].font.size = Pt(10)
    return prs

def word2pic(wordPath):
    """
    extract picture from HPLC word file
    return the priture name list
    """
    picList = []
    doc = docx.Document(wordPath)
    proNum = len(doc.tables)
    for i in range(proNum):
        for row in range(5):
            for cell in range(5):
                if doc.tables[i].rows[row].cells[cell].text == 'Sample name:':
                    sampleName = doc.tables[i].rows[row].cells[cell + 1].text
        blip = doc.inline_shapes[2*i]._inline.graphic.graphicData.pic.blipFill.blip
        rID = blip.embed
        image_part = doc.part.related_parts[rID]
        with open(sampleName+'.bin', "wb") as f:
            f.write(image_part.blob)
        picList.append(sampleName+'.bin')
    return picList

def pdf2pic(pdfPath):
    """
    extract picture from HPLC pdf file
    return the priture name list
    """
    picList = []
    fp = open(pdfPath,'rb')
    with open(pdfPath,'rb') as fp:
        praser = PDFParser(fp)
        doc = PDFDocument()
        praser.set_document(doc)
        doc.set_parser(praser)
        doc.initialize()
        if not doc.is_extractable:
            raise PDFTextExtractionNotAllowed
        else:
            rsrcmgr = PDFResourceManager()
            laparams = LAParams()
            device = PDFPageAggregator(rsrcmgr, laparams=laparams)
            interpreter = PDFPageInterpreter(rsrcmgr, device)
            results = []
            for page in doc.get_pages():
                interpreter.process_page(page)
                layout = device.get_result()
                for x in layout:
                    if isinstance(x, LTTextBoxHorizontal):
                        results.append(x.get_text())
    sampleName = [results[i*18 + 2].split('\n')[0] for i in range(int(len(results)/18))]
    
    checkXO = r"/Type(?= */XObject)"
    checkIM = r"/Subtype(?= */Image)"
    doc = fitz.open(pdfPath)
    lenXREF = doc._getXrefLength()
    picOrder = 0
    for i in range(1, lenXREF):
        text = doc._getXrefString(i)
        isXObject = re.search(checkXO, text)
        isImage = re.search(checkIM, text)
        if isXObject and isImage:
            picOrder += 1
            if picOrder % 2 == 0:
                pix = fitz.Pixmap(doc, i)
                picName = sampleName[int(picOrder // 2)-1] + '.png'
                picList.append(picName)
                if pix.n < 5:
                    pix.writePNG(picName)
                else:
                    pix0 = fitz.Pixmap(fitz.csRGB, pix)
                    pix0.writePNG(picName)
                    pix0 = None
                pix = None
    return picList

if __name__ == '__main__':
    filePath = os.path.abspath(os.path.dirname(__file__))
    template_path = os.path.join(filePath,'purificationTemplate.pptx')
    prs = Presentation(template_path)
    hplcList = [
        ['1','Protein A','Protein A.docx','WXBIOXXXX/XXX/20200102',
        "The purities of the proteins were below 90%.|The retention time \
            around at 7.7 mins indicating the proteins monomers."],
        ['2','Dialysis','Dialysis.pdf','WXBIOXXX/XXX/20200110',
        "The purities of the proteins were above 90%.|The retention time \
        around at 7.7 mins indicating the proteins monomers."]
        ] 
    proDict ={            
        '1_Protein A':{
            'W308031-T6U9.E17-1.uIgG1': [1,'59.77%'],
            'W308031-U9T6.E17-1.uIgG1': [2,'86.79%'],
            'W308031-T6U9.E17R-1.uIgG1': [3,'85.68%'],
            'W308031-U9T6.E17R-1.uIgG1': [4,'75.10%']
            },
        '2_Dialysis':{
            'W308031-T6U9.E17-1.uIgG1': [1,'99.62%'],
            'W308031-T6U9.E17R-1.uIgG1': [3,'99.30%']
            }
            }
    prs = purificationHplcPageMake(prs,hplcList,proDict)
    prs.save('purificationHplcPageTest.pptx')
