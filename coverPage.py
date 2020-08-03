# -*- coding: utf-8 -*-
"""
Created on Wed Jun 10 16:46:15 2020
@author: tao_jingfen
"""
from pptx import Presentation
import os

def purificationCoverPageMake(prs: str, projectName: str, date: str) -> str:
    """
    Create the cover page of purification auto-report PPT

    --Args:
        prs: presentation object loaded from *pptx* by template pptx
        projectName：the name of the project, e.g. WBP123456 
        date: the date of report, format like "04/13/2020"
    --Returns:
        prs: presentation object could be used to save pptx
    """
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    shapes = slide.shapes
    shapes.placeholders[0].text =  projectName + ' Purification Report'
    shapes.placeholders[1].text = 'BID_PE'
    date_paragraph = shapes.placeholders[1].text_frame.add_paragraph()
    date_paragraph.text = date
    return prs

if __name__ == '__main__':
    filePath = os.path.abspath(os.path.dirname(__file__))
    template_path = os.path.join(filePath,'purificationTemplate.pptx')
    prs = Presentation(template_path)
    projectName = 'WBPXXX'
    date = '04/13/2020'
    prs = purificationCoverPageMake(prs, projectName, date)
    prs.save('purificationCoverPageTest.pptx')